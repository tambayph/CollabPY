from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime, timedelta
import os, sys


# base path
base_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'data')
folder = os.path.dirname(os.path.abspath(__file__))
script_name = sys.argv[1]

# Function to add a slide with or without an image
def add_slide_with_image(presentation, title, date, image_path, layout_index):
    slide_layout = presentation.slide_layouts[layout_index]
    slide = presentation.slides.add_slide(slide_layout)

   # Add title at the center horizontally and with a top margin, and make it bold with font size 18pt
    title_shape = slide.shapes.title
    if title_shape:
        # Concatenate "yesterday" date with the title
        title_text = f"{title} - {date}"
        title_shape.text = title_text
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.size = Pt(24)

        # Set word wrap to True for horizontal text
        title_shape.text_frame.word_wrap = True

        # Calculate the center position for the title horizontally and set width to 50% of slide width
        title_width_percent = 100
        title_width = int((presentation.slide_width * title_width_percent) / 100)
        title_left = int((presentation.slide_width - title_width) / 2)

        # Set the top margin
        top_margin = Pt(20)
        title_top = int(top_margin)

        # Assign the calculated values
        title_shape.left = title_left
        title_shape.top = title_top
        title_shape.width = title_width

        # Explicitly set the text orientation to horizontal
        title_shape.text_frame.text_orientation = 0  # 0 degrees for horizontal orientation

    # Add image to the slide
    if image_path:
        # Calculate the center position for the main image
        main_width = Inches(10)  # Adjust the width as needed
        main_height = Inches(6)
        main_left = (presentation.slide_width - main_width) / 10  # Centered within the first half
        main_top = (presentation.slide_height - main_height) / 2

        slide.shapes.add_picture(image_path, main_left, main_top, width=main_width, height=main_height)

    # Add the date at the bottom and make it bold
    placeholders_count = len(slide.placeholders)
    subheader_shape = None

    # Check if there is a placeholder at index 1
    if placeholders_count > 1:
        subheader_shape = slide.placeholders[1].text_frame.add_paragraph()
        subheader_shape.text = date
        subheader_shape.font.bold = True
    else:
        print(f"Warning: No placeholder at index 1 in layout {layout_index}. Subheader not added.")


# Main function to create the presentation
def create_presentation(selected_date):
     # Parse the input string as a datetime object
    datetime_object = datetime.strptime(selected_date, '%Y%m%d')

    # Format the date in 'dd Month YYYY' format
    formatted_date = datetime_object.strftime('%d %B %Y')

  
    # Create a PowerPoint presentation
    presentation = Presentation()

    # Slide 1: Map Discussion (Title Slide)
    title = "MAP DISCUSSION"
    current_date = datetime_object
    add_slide_with_image(presentation, title, formatted_date, None, 0)  

    # Ensure current_date is not a string; reassign only if needed
    if not isinstance(current_date, datetime):
        current_date = datetime.strptime(current_date, '%d %B %Y')

    # Subtract one day from the current date
    yesterday_date = current_date - timedelta(days=1)
    yesterday_date_str = yesterday_date.strftime('%d %B %Y')


    # Slide 2: Surface Chart, 00z
    title = "Surface Chart, 00Z"

    image_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'data', 'surface_map.jpg')
    #  upper air code
    levels = [850, 700, 500, 300, 200]

    # Use a loop to create paths for each level
    image_paths = [os.path.join(base_path, f'{level}MB.jpg') for level in levels]

    # Now image_paths contains the paths for each level
    image_850_path, image_700_path, image_500_path, image_300_path, image_200_path = image_paths

    add_slide_with_image(presentation, title, yesterday_date_str, image_path, 5)  # Use layout index 5 for content slide

    # Slide 3: Surface Chart, 03z
    title = "Surface Chart, 06Z"
    add_slide_with_image(presentation, title, yesterday_date_str, image_path, 5)  # Use layout index 5 for content slide

    # Slide 4: Surface Chart, 012z
    title = "Surface Chart, 12Z"
    add_slide_with_image(presentation, title, yesterday_date_str, image_path, 5)  # Use layout index 5 for content slide

    # Slide 5: Surface Chart, 18Z
    title = "Surface Chart, 18Z"
    add_slide_with_image(presentation, title, yesterday_date_str, image_path, 5)  # Use layout index 5 for content slide

    # Slide 6: Surface Chart, 18Z
    title = "Surface Chart, 00Z"
    add_slide_with_image(presentation, title, formatted_date, image_path, 5)  # Use layout index 5 for content slide

    # Slide 6: Surface Chart, 00Z
    title = "(TMD) Surface Chart, 00Z"
    add_slide_with_image(presentation, title, formatted_date, image_path, 5)  # Use layout index 5 for content slide

    # Slide 7: Surface Chart, 18Z
    title = "(JMA)Surface Chart, 00Z"
    add_slide_with_image(presentation, title, formatted_date, image_path, 5)  # Use layout index 5 for content slide

    # Slide 8: Upper air chart, 
    title = "Upper Air Chart, 850MB, 00Z"
    add_slide_with_image(presentation, title, formatted_date, image_850_path, 5)  # Use layout index 5 for content slide

    # Slide 9: Upper air chart, 
    title = "Upper Air Chart, 700MB, 00Z"
    add_slide_with_image(presentation, title, formatted_date, image_700_path, 5)  # Use layout index 5 for content slide

    # Slide 10: Upper air chart, 
    title = "Upper Air Chart, 500MB, 00Z"
    add_slide_with_image(presentation, title, formatted_date, image_500_path, 5)  # Use layout index 5 for content slide

    # Slide 11: Upper air chart, 
    title = "Upper Air Chart, 300MB, 00Z"
    add_slide_with_image(presentation, title, formatted_date, image_300_path, 5)  # Use layout index 5 for content slide

    # Slide 12: Upper air chart, 
    title = "Upper Air Chart, 200MB, 00Z"
    add_slide_with_image(presentation, title, formatted_date, image_200_path, 5)  # Use layout index 5 for content slide

    # Slide 13: Suggested Scenario 
    title = "Suggested Scenario"
    add_slide_with_image(presentation, title, formatted_date, '',5)  # Use layout index 5 for content slide


   # Save the single PowerPoint presentation to a file
    presentation.save(f'{folder}/powerpoints/MapDiscussion_{script_name}.pptx')

if __name__ == "__main__":

    # Check if the script is called with the expected number of arguments
    if len(sys.argv) != 2:
        print("Usage: python Map_Discussion.py <datetime_string>")
        sys.exit(1)

    # Extract the datetime_string from the command line arguments
    datetime_string = sys.argv[1]

    create_presentation(datetime_string)