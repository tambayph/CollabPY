import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os,sys
from datetime import datetime,timedelta
from pptx.enum.shapes import MSO_SHAPE_TYPE


####################################
#declaring absolute path
folder = os.path.dirname(os.path.abspath(__file__))
#for api datetime
script_name = sys.argv[1]
###################################

# Function to add legend image to slide
def add_legend_image(slide, image_path, left, top, width, height):
    if os.path.exists(image_path):
        # Add the image to the slide
        slide.shapes.add_picture(image_path, left, top, width=width, height=height)
    else:
        print(f"Image file not found: {image_path}")

# Fetch data from the API
rainfall_url = f'http://10.11.1.107/api/rainfall/{script_name}'
rainfall_response = requests.get(rainfall_url)

cloud_cover_url = f'http://10.11.1.107/api/cloud-cover/{script_name}'
cloud_cover_response = requests.get(cloud_cover_url)

if rainfall_response.status_code == 200 and cloud_cover_response.status_code == 200:
    rainfall_data = rainfall_response.json()
    cloud_cover_data = cloud_cover_response.json()
else:
    print("Failed to fetch data from the API")
    exit()

# Create PowerPoint presentation
prs = Presentation()

# Define font properties
font_name = 'Arial'
font_size = Pt(12)

# Calculate the number of modified data columns
num_modified_columns = len(rainfall_data[0]["rainfall"])

# Define table dimensions and position
top = Inches(0)
left = Inches(0)
width = Inches(7.8)  # Adjusted to fit the whole slide without any margin
height = Inches(6)  # Adjusted to fit the whole slide without any margin

# Add content slides with tables
for island_category in ['Luzon', 'Visayas', 'Mindanao']:
    # Filter stations by island category
    filtered_stations_rainfall = [station for station in rainfall_data if station.get('island_category', '').lower() == island_category.lower()]
    filtered_stations_cloud_cover = [station for station in cloud_cover_data if station.get('island_category', '').lower() == island_category.lower()]

    # Skip empty categories
    if not filtered_stations_rainfall or not filtered_stations_cloud_cover:
        continue

    # Calculate number of slides needed for this category
    num_rows_per_slide_category = 10  # Each slide for Visayas and Mindanao will have 10 rows
    num_slides_category = (max(len(filtered_stations_rainfall), len(filtered_stations_cloud_cover)) + num_rows_per_slide_category - 1) // num_rows_per_slide_category

    print(f"DEBUG: Category: {island_category}, Num slides: {num_slides_category}")

    for slide_index in range(num_slides_category):
        num_rows = num_rows_per_slide_category  # Each slide for Visayas and Mindanao will have 10 rows
        
        print(f"DEBUG: Slide index: {slide_index}, Num rows: {num_rows}")

        slide_layout = prs.slide_layouts[5]  # Blank slide layout with title and content
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        table = shapes.add_table(rows=(num_rows + 1) * 2, cols=num_modified_columns + 2, left=left, top=top, width=width, height=height).table

        # Set font properties for headers and data cells
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = font_name
                        run.font.size = font_size

        # Add headers for rainfall
        table.cell(0, 0).text = 'Luzon' if island_category.lower() == 'luzon' else 'Visayas' if island_category.lower() == 'visayas' else 'Mindanao'

        for i, modified_data in enumerate(rainfall_data[0]['rainfall'], start=1):
            header_text = modified_data['dateTimeUTC'][11:13] + 'Z'
            table.cell(0, i).text = header_text
        
        # Add 24hr column for rainfall
        table.cell(0, num_modified_columns + 1).text = "24hr"

        # Add headers for cloud cover
        for i, cloud_cover in enumerate(cloud_cover_data[0]['cloud_cover'], start=1):
            header_text = cloud_cover['dateTimeUTC'][11:13] + 'Z'
            table.cell(num_rows + 1, i).text = header_text

        # Add 24hr column for cloud cover
        table.cell(num_rows + 1, num_modified_columns + 1).text = "24hr"

        # Set font properties for headers and data cells
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(13)

        # Set font size for the 24hr column headers
        table.cell(0, num_modified_columns + 1).text_frame.paragraphs[0].runs[0].font.size = Pt(13)
        table.cell(num_rows + 1, num_modified_columns + 1).text_frame.paragraphs[0].runs[0].font.size = Pt(13)


        # Add station info and modified data
        prev_station_name = None

        for i in range(num_rows):
            station_index = slide_index * num_rows + i
            if station_index >= max(len(filtered_stations_rainfall), len(filtered_stations_cloud_cover)):
                break

            # Add rainfall data
            if station_index < len(filtered_stations_rainfall):
                station = filtered_stations_rainfall[station_index]
                station_name = f"{station['stn_name']}"
                table.cell(i * 2 + 1, 0).text = station_name

                for j, modified_data in enumerate(station['rainfall'], start=1):
                    value = modified_data['value']
                    table.cell(i * 2 + 1, j).text = value
                    # Limit the number of rows per cell to 20
                    table.cell(i * 2 + 1, j).text_frame.paragraphs[0].runs[0].font.size = Pt(12)
                    table.cell(i * 2 + 1, j).text_frame.paragraphs[0].space_after = Pt(12)
                    table.cell(i * 2 + 1, j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                # Calculate the difference between the values in cell (i * 2 + 1, 2) and (i * 2 + 1, 1)
                value_1_text = table.cell(i * 2 + 1, 2).text
                value_2_text = table.cell(i * 2 + 1, 1).text

                if value_1_text == 'T':
                    value_1 = 'T'  # Set value_1 to 'T'
                elif value_1_text == '':
                    value_1 = 0.0  # Assign 0.0 if the string is empty
                else:
                    value_1 = float(value_1_text)

                if value_2_text == 'T':
                    value_2 = 'T'  # Set value_2 to 'T'
                elif value_2_text == '':
                    value_2 = 0.0  # Assign 0.0 if the string is empty
                else:
                    value_2 = float(value_2_text)

                # Calculate difference only if both values are not 'T'
                if value_1 != 'T' and value_2 != 'T':
                    difference = value_1 - value_2
                else:
                    difference = 'T'  # Set difference to 'T' if either value is 'T'

                # Update the cell (i * 2 + 1, 2) with the difference including the decimal part
                if difference != 'T':
                    table.cell(i * 2 + 1, 2).text = f"{difference:.1f}"
                else:
                    table.cell(i * 2 + 1, 2).text = value_1_text

                # Set font size to 12pt for the updated cell
                table.cell(i * 2 + 1, 2).text_frame.paragraphs[0].runs[0].font.size = Pt(12)

                                 # Calculate the difference between the values in cell (i * 2 + 1, 4) and (i * 2 + 1, 3)
                value_3_text = table.cell(i * 2 + 1, 4).text
                value_4_text = table.cell(i * 2 + 1, 3).text

                if value_3_text == 'T':
                    value_3 = 'T'  # Set value_3 to 'T'
                elif value_3_text == '':
                    value_3 = 0.0  # Assign 0.0 if the string is empty
                else:
                    value_3 = float(value_3_text)

                if value_4_text == 'T':
                    value_4 = 'T'  # Set value_4 to 'T'
                elif value_4_text == '':
                    value_4 = 0.0  # Assign 0.0 if the string is empty
                else:
                    value_4 = float(value_4_text)

                # Calculate difference only if both values are not 'T'
                if value_3 != 'T' and value_4 != 'T':
                    difference = value_3 - value_4
                else:
                    difference = 'T'  # Set difference to 'T' if either value is 'T'

                # Update the cell (i * 2 + 1, 4) with the difference including the decimal part
                if difference != 'T':
                    table.cell(i * 2 + 1, 4).text = f"{difference:.1f}"
                else:
                    table.cell(i * 2 + 1, 4).text = value_3_text

                # Set font size to 12pt for the updated cell
                table.cell(i * 2 + 1, 4).text_frame.paragraphs[0].runs[0].font.size = Pt(12)

                 # Calculate the difference between the values in cell (i * 2 + 1, 6) and (i * 2 + 1, 5)
                value_5_text = table.cell(i * 2 + 1, 6).text
                value_6_text = table.cell(i * 2 + 1, 5).text

                if value_5_text == 'T':
                    value_5 = 'T'  # Set value_5 to 'T'
                elif value_5_text == '':
                    value_5 = 0.0  # Assign 0.0 if the string is empty
                else:
                    value_5 = float(value_5_text)

                if value_6_text == 'T':
                    value_6 = 'T'  # Set value_6 to 'T'
                elif value_6_text == '':
                    value_6 = 0.0  # Assign 0.0 if the string is empty
                else:
                    value_6 = float(value_6_text)

                # Calculate difference only if both values are not 'T'
                if value_5 != 'T' and value_6 != 'T':
                    difference = value_5 - value_6
                else:
                    difference = 'T'  # Set difference to 'T' if either value is 'T'

                # Update the cell (i * 2 + 1, 6) with the difference including the decimal part
                if difference != 'T':
                    table.cell(i * 2 + 1, 6).text = f"{difference:.1f}"
                else:
                    table.cell(i * 2 + 1, 6).text = value_5_text

                # Set font size to 12pt for the updated cell
                table.cell(i * 2 + 1, 6).text_frame.paragraphs[0].runs[0].font.size = Pt(12)

                # Calculate the difference between the values in cell (i * 2 + 1, 8) and (i * 2 + 1, 7)
                value_7_text = table.cell(i * 2 + 1, 8).text
                value_8_text = table.cell(i * 2 + 1, 7).text

                if value_7_text == 'T':
                    value_7 = 'T'  # Set value_7 to 'T'
                elif value_7_text == '':
                    value_7 = 0.0  # Assign 0.0 if the string is empty
                else:
                    value_7 = float(value_7_text)

                if value_8_text == 'T':
                    value_8 = 'T'  # Set value_8 to 'T'
                elif value_8_text == '':
                    value_8 = 0.0  # Assign 0.0 if the string is empty
                else:
                    value_8 = float(value_8_text)

                # Calculate difference only if both values are not 'T'
                if value_7 != 'T' and value_8 != 'T':
                    difference = abs(value_7 - value_8)
                else:
                    difference = 'T'  # Set difference to 'T' if either value is 'T'

                # Update the cell (i * 2 + 1, 8) with the difference including the decimal part
                if difference != 'T':
                    table.cell(i * 2 + 1, 8).text = f"{difference:.1f}"
                else:
                    table.cell(i * 2 + 1, 8).text = value_7_text

                # Set font size to 12pt for the updated cell
                table.cell(i * 2 + 1, 8).text_frame.paragraphs[0].runs[0].font.size = Pt(12)


               


                # Apply fill color green to cells from column 1 to 8
                for col_idx in range(1, 9):  # Columns 1 to 8
                    #Iterate through rows (assuming 'i' represents row index)
                    for row_idx in range(1, num_rows + 1):
                        try:
                            cell_value = float(table.cell(i * 2 + 1, col_idx).text)
                            fill_color = RGBColor(255, 255, 255)  # Default to green
                            
                            # Apply condition for yellow fill color when cell value < 7.5
                            if cell_value  == 0:
                                fill_color = RGBColor(255, 255, 255)  # white
                            elif cell_value < 7.5:
                                fill_color = RGBColor(255, 255, 0)  # Yellow
                            elif 7.5 <= cell_value < 22.5:
                                fill_color = RGBColor(255, 165, 0)  # Orange
                            elif 22.5 <= cell_value < 45:
                                fill_color = RGBColor(255, 0, 0)  # Red
                            elif 45 <= cell_value < 90:
                                fill_color = RGBColor(255, 192, 203)  # Pink
                            # Apply condition for violet fill color when cell value is greater than 90
                            elif cell_value >= 90:
                                fill_color = RGBColor(238, 130, 238)  # Violet
                            
                            fill = table.cell(i * 2 + 1, col_idx).fill
                            fill.solid()
                            fill.fore_color.rgb = fill_color
                        except ValueError:
                            pass  # If the cell doesn't contain a float, skip coloring

                    
                # Add 24hr data for rainfall from JSON
                if 'total' in station and station['total'] == "T":
                    table.cell(i * 2 + 1, num_modified_columns + 1).text = "T"
                    table.cell(i * 2 + 1, num_modified_columns + 1).text_frame.paragraphs[0].runs[0].font.size = Pt(13)
                else:
                    if 'total' in station and len(station['total']) > 0:
                        total_value = float(station['total'])
                        formatted_total = "{:.1f}".format(total_value)
                        table.cell(i * 2 + 1, num_modified_columns + 1).text = formatted_total
                        table.cell(i * 2 + 1, num_modified_columns + 1).text_frame.paragraphs[0].runs[0].font.size = Pt(13)
                    else:
                        table.cell(i * 2 + 1, num_modified_columns + 1).text = "0.0"
                        table.cell(i * 2 + 1, num_modified_columns + 1).text_frame.paragraphs[0].runs[0].font.size = Pt(13)


            # Add cloud cover data
            if station_index < len(filtered_stations_cloud_cover):
                station = filtered_stations_cloud_cover[station_index]
                prev_station_name = f"{station['stn_name']}"
                #  table.cell(i * 2 + 2, 0).text = station_name
                
                # Merge cells vertically for station name
                if station_name == prev_station_name:
                    if not table.cell(i * 2 + 1, 0).is_spanned:
                        table.cell(i * 2 + 2, 0).merge(table.cell(i * 2 + 1, 0))
                else:
                    prev_station_name = station_name

                for j, cloud_cover in enumerate(station['cloud_cover'], start=1):
                    if cloud_cover['oktas'] == '':
                        value_text = "-"
                    else:
                        value_text = f"{cloud_cover['oktas']}/{cloud_cover['cloudgroup']}"
                    table.cell(i * 2 + 2, j).text = value_text
                    # Limit the number of rows per cell to 20
                    table.cell(i * 2 + 2, j).text_frame.paragraphs[0].runs[0].font.size = Pt(12)
                    table.cell(i * 2 + 2, j).text_frame.paragraphs[0].space_after = Pt(12)
                    table.cell(i * 2 + 1, j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
               
                # Merge the last row vertically with the previous row in the last column
                table.cell(i * 2 + 2, num_modified_columns + 1).merge(table.cell(i * 2 + 1, num_modified_columns + 1))
              
                # Center the text within the merged cell
                table.cell(i * 2 + 2, num_modified_columns + 1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
          
        # Set font properties for the first column (station number & name)
        first_column_cells = [table.cell(i, 0) for i in range((num_rows + 1) * 2)]
        for cell in first_column_cells:
            cell.text_frame.paragraphs[0].font.name = font_name
            cell.text_frame.paragraphs[0].font.size = font_size

        # Set width for the first column
        table.columns[0].width = Inches(3)  # Adjusted to fit the whole slide without any margin

# Get the index of the last column
last_column_index = len(table.columns) - 1

# Iterate over all slides in the presentation
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        table = shape.table
        
        #  add - to empty cell 
        # Iterate over all cells in the table
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                # Check if the cell text is empty
                if not cell.text.strip():
                    # Replace empty cell with "-"
                    cell.text = "-"
                    
                # Set alignment to center for each cell
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
        # Iterate over table rows starting from the second row
        for idx, row in enumerate(table.rows):
            if idx == 0:
                continue  # Skip the first row

            # Get the value of the last column in the row
            cell_value = row.cells[last_column_index].text

            # Convert the cell value to float if it's not '-' (assuming it's numeric data)
            if cell_value.strip():  # Check if the cell value is not empty or blank
                try:
                    column_value = float(cell_value)
                except ValueError:
                    column_value = 0.0  # Set to 0 if conversion fails
            else:
                # If the cell value is empty or blank, assign a specific color and set column_value to 0
                column_value = 0.0
                fill_color = RGBColor(255, 255, 255)  # White for blank or empty cells

            # Apply fill color based on the value
            if column_value == 0.0:
                fill_color = RGBColor(169, 169, 169)  # Gray
            elif 0 <= column_value <= 60:
                fill_color = RGBColor(255, 255, 0)  # Yellow
            elif 60 < column_value <= 180:
                fill_color = RGBColor(255, 165, 0)  # Orange
            elif 180 < column_value <= 360:
                fill_color = RGBColor(255, 0, 0)  # Red
            elif 360 < column_value <= 720:
                fill_color = RGBColor(255, 192, 203)  # Pink
            elif column_value >= 720:
                fill_color = RGBColor(238, 130, 238)  # Violet
            else:
                fill_color = RGBColor(255, 255, 255)  # White (default)
            
            # Set the cell fill color
            fill = row.cells[last_column_index].fill
            fill.solid()
            fill.fore_color.rgb = fill_color

            # Set text alignment to center
            row.cells[last_column_index].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Check if the table has any rows

            if len(table.rows) > 0:
                last_row_idx = len(table.rows) - 1
                last_row = table.rows[last_row_idx]
                background_color = RGBColor(255, 255, 255) 
                for cell in last_row.cells:
                    cell.text = ""
                    fill = cell.fill
                    fill.solid()
                    fill.fore_color.rgb = background_color


            # for 4th slide
            for idx, slide in enumerate(prs.slides):
                if idx == 3:  # Check if it's the fourth slide (index starts from 0)
                    for shape in slide.shapes:
                        if not shape.has_table:
                            continue
                        table = shape.table

                        # Check if the table has at least three rows
                        if len(table.rows) > 4:
                            # Iterate over the last three rows of the table
                            for row_idx in range(len(table.rows) - 5, len(table.rows)):
                                row = table.rows[row_idx]
                                # Set the background color of each cell in the row to white
                                for cell in row.cells:
                                    cell.text = ""
                                    fill = cell.fill
                                    fill.solid()
                                    fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
            # make last row with gray into white
            # Iterate over all slides
            for slide in prs.slides:
                # Iterate over all shapes in the slide
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        # Iterate over all rows in the table
                        for row in shape.table.rows:
                            # Iterate over all cells in the row
                            for cell in row.cells:
                                # Check if the cell is blank
                                if not cell.text.strip():
                                    # Set the background color of the cell to white
                                    fill = cell.fill
                                    fill.solid()
                                    fill.fore_color.rgb = RGBColor(255, 255, 255)  # White

# Define legend URLs
folder = os.path.dirname(os.path.abspath(__file__))
rainfall_legend_url = f'{folder}/images/rainfall_legend.png'
cloud_legend_url = f'{folder}/images/cloud_legend.png'

# Define legend coordinates and dimensions
rainfall_legend_left = Inches(5)  # Adjust as needed
rainfall_legend_top = Inches(6.4)  # Adjust as needed
rainfall_legend_width = Inches(5.0)  # Adjust as needed
rainfall_legend_height = Inches(1.1)  # Adjust as needed
rainfall_legend_right = Inches(0)  # Adjust as needed

cloud_legend_left = Inches(0)  # Adjust as needed
cloud_legend_top = Inches(6.4)  # Adjust as needed
cloud_legend_width = Inches(5.1)  # Adjust as needed
cloud_legend_height = Inches(1.1)  # Adjust as needed

# Add legends to each slide
for slide in prs.slides:
    # Add rainfall legend
    add_legend_image(slide, rainfall_legend_url, rainfall_legend_left, rainfall_legend_top, rainfall_legend_width, rainfall_legend_height)
    
    # Add cloud cover legend
    add_legend_image(slide, cloud_legend_url, cloud_legend_left, cloud_legend_top, cloud_legend_width, cloud_legend_height)

# Add legends to each slide
for slide in prs.slides:
    # Add rainfall legend
    add_legend_image(slide, rainfall_legend_url, rainfall_legend_left, rainfall_legend_top, rainfall_legend_width, rainfall_legend_height)

    # Add cloud cover legend
    add_legend_image(slide, cloud_legend_url, cloud_legend_left, cloud_legend_top, cloud_legend_width, cloud_legend_height)

# Save the single PowerPoint presentation to a file

script_name = sys.argv[1]
prs.save(f'{folder}/powerpoints/MapDiscussion_{script_name}.pptx')


# COPY TO TERA CODE
transfer = os.path.join(folder, '../copytotera')
sys.path.append(transfer)

from copytotera import mount_network_server, transfer_to_nas, nas_local_dir

# Define the base remote path
base_remote_path = 'wd.s.dstor.pagasa.local/wfs/Rainfall_table_for_Map_Discussion/'

# Get current year and month
current_year = datetime.now().strftime("%Y")
current_year_2_digit = datetime.now().strftime("%y")
current_month = datetime.now().strftime("%m")
month_name = datetime.now().strftime("%b").upper()

# Convert the current month to 3-character format with leading zeros
month_number = datetime.now().strftime("%m")

# Construct the remote path with current year and month
nas_remote_path = f"{base_remote_path}{current_year}/{month_number} {month_name}'{current_year_2_digit}"

# Mount the network server
mount_network_server(nas_remote_path, nas_local_dir)

# Transfer the source image to the NAS local directory
transfer_to_nas(f'{folder}/powerpoints/MapDiscussion_{script_name}.pptx', nas_local_dir)
