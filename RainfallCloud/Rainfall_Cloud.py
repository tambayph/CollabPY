from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from datetime import datetime,timedelta
import requests
import os,sys
from views import get_absolute_image_path




####################################
#declaring absolute path
folder = os.path.dirname(os.path.abspath(__file__))
#for api datetime
script_name = sys.argv[1]
###################################
def add_legend_image(slide, image_path, left, top, width, height):
    slide.shapes.add_picture(image_path, left, top, width, height)


def add_default_slide(presentation, main_image_filename, additional_image_filename):
    # Create a new slide with a blank layout
    slide_layout = presentation.slide_layouts[6]  # Assuming layout index 6 is blank
    slide = presentation.slides.add_slide(slide_layout)

    # Use get_absolute_image_path to get the absolute path to the main image
    absolute_main_image_path = get_absolute_image_path(main_image_filename)

    # Calculate the center position for the main image
    main_width = Inches(4)  # Adjust the width as needed
    main_height = Inches(6)
    main_left = (presentation.slide_width - main_width) / 10  # Centered within the first half
    main_top = (presentation.slide_height - main_height) / 2

    # Add the main image to the slide using the absolute path
    slide.shapes.add_picture(absolute_main_image_path, main_left, main_top, main_width, main_height)

    # Use get_absolute_image_path to get the absolute path to the additional image
    absolute_additional_image_path = get_absolute_image_path(additional_image_filename)

    # Calculate the center position for the additional image
    additional_width = Inches(4)  # Adjust the width as needed
    additional_height = Inches(6)
    additional_left = (presentation.slide_width - additional_width) * 3 / 4  # Centered within the second half
    additional_top = main_top

    # Add the additional image to the first slide using the absolute path
    slide.shapes.add_picture(absolute_additional_image_path, additional_left, additional_top, additional_width, additional_height)


def create_presentation(presentation, island_category, rainfall_json_file_url, cloud_json_file_url):
    # Retrieve JSON data from the rainfall URL
    response = requests.get(rainfall_json_file_url)
    if response.status_code == 200:
        rainfall_data = response.json()
    else:
        print(f"Failed to retrieve JSON data from {rainfall_json_file_url}. Status code: {response.status_code}")
        return

    # Retrieve JSON data from the cloud cover URL
    response = requests.get(cloud_json_file_url)
    if response.status_code == 200:
        cloud_cover_data = response.json()
    else:
        print(f"Failed to retrieve JSON data from {cloud_json_file_url}. Status code: {response.status_code}")
        return

    # Filter data for the specified island category
    island_rainfall_data = [station for station in rainfall_data if station.get("island_category") == island_category]
   

    # Define the maximum number of rows per slide
    max_rows_per_slide = 20
    
    # Extract datetime UTC values and rainfall values from the first station
    datetime_utc_values = []
    for entry in island_rainfall_data[0]["rainfall"]:
        if isinstance(entry, dict) and "dateTimeUTC" in entry:
            datetime_utc_values.append(entry["dateTimeUTC"])
    rainfall_values = []
    for entry in island_rainfall_data[0]["rainfall"]:
        if isinstance(entry, dict) and "value" in entry:
            value = entry["value"]
            try:
                value = float(value)  # Attempt to convert to float
            except ValueError:
                # Handle non-numeric values (e.g., replace 'T' with 0.0)
                value = 0.0
            rainfall_values.append(value)

    # Create a set to store all unique dateTimeUTC values
    unique_datetime_utc_values = []
    seen_datetimes = set()

    # Iterate through the data and add unique datetime values to the set
    for station_data in island_rainfall_data:
        for entry in station_data["rainfall"]:
            if isinstance(entry, dict):
                try:
                    value = float(entry["value"])
                except ValueError:
                    # Handle non-numeric values (e.g., replace 'T' with 0.0)
                    value = 0.0

                datetime_utc = entry["dateTimeUTC"]
                if datetime_utc not in seen_datetimes:
                    unique_datetime_utc_values.append(datetime_utc)
                    seen_datetimes.add(datetime_utc)

    # Convert the set to a list (no sorting)
    unique_datetime_utc_values = list(unique_datetime_utc_values)

    # print(unique_datetime_utc_values)
    # Define the header row data (including "12z Diff" column) without duplicates
    header_row = [island_category] + [f"{datetime_obj.strftime('%Hz')}" for datetime_obj in
                                    map(lambda dt: datetime.strptime(dt, "%Y-%m-%dT%H:%M:%S.%fZ"),
                                        unique_datetime_utc_values)] + ["24HR"]

    # Calculate the number of columns for the table
    
    num_cols = len(datetime_utc_values) 
    expected_num_cols = len(datetime_utc_values) + 1  # Add columns for differences and "00z"
    expected_num_cols_greater = len(datetime_utc_values) + 2  # Add columns for differences and "00z"

    if num_cols > expected_num_cols:
        num_cols = expected_num_cols
    else:
        num_cols = expected_num_cols_greater
            
    # Initialize row index
    row_idx = 1
    
    # Define font family and font size for all cells in the table
    font_family = "Arial"
    font_size = Pt(10)

    # Keep track of the previous station name
    prev_station_name = None

    # Filter stations for display only if the island_category is the specified one
    island_stations = [station for station in rainfall_data if station.get("island_category") == island_category]
    island_cloud_stations = [station for station in cloud_cover_data if station.get("island_category") == island_category]
    
    # Initialize a flag for alternating between rainfall and cloud data
    is_rainfall_data = True

    # Iterate through the filtered data
    for station_data, cloud_station_data in zip(island_stations, island_cloud_stations):
        # Determine the data type based on the flag
        if is_rainfall_data:
            data_type = "rainfall"  # Use "rainfall" data type
            data_entries = station_data[data_type]
        else:
            data_type = "cloud_cover"  # Use "cloud_cover" data type
            data_entries = cloud_station_data[data_type]
        for data_type in ["rainfall", "cloud_cover"]:
            # Create a new slide for the table (if needed)
            if row_idx == 1:
                current_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
                table = current_slide.shapes.add_table(
                    max_rows_per_slide + 1, num_cols, Inches(0), Inches(0), Inches(0), Inches(0)
                ).table

                # Set fixed widths for columns on the new slide
                table.columns[0].width = Inches(1.9)  # Station Name
                for col_idx in range(1, num_cols - 1):
                    table.columns[col_idx].width = Inches(.9)  # Other datetime columns
                table.columns[num_cols - 1].width = Inches(.9)  # "00z" column

                # Add headers to the first row of the table
                for col_idx, header_text in enumerate(header_row):
                    table.cell(0, col_idx).text = header_text

                # Set font size and alignment for header cells
                for col_idx in range(num_cols):
                    cell = table.cell(0, col_idx)
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = font_size
                            run.font.name = font_family  # Set the font family here
                            paragraph.alignment = PP_ALIGN.CENTER

                # Set background color and border for header cells
                header_cell = table.cell(0, 0)  # Assuming the first cell is the header cell
                header_cell.fill.solid()
                header_cell.fill.fore_color.rgb = RGBColor(91, 155, 213)  # Blue color for the header
                for cell in header_row[1:]:
                    header_table_cell = table.cell(0, header_row.index(cell))
                    header_table_cell.fill.solid()
                    header_table_cell.fill.fore_color.rgb = RGBColor(91, 155, 213)  # Blue color for the header

            # Select the appropriate data based on the data type ("rainfall" or "cloud_cover")
            if data_type == "rainfall":
                data_entries = station_data[data_type]
                station_name = station_data["stn_name"].split(',')[0]
            else:
                data_entries = cloud_station_data[data_type]
                station_name = station_data["stn_name"].split(',')[0]

 
            # If the station name is the same as the previous one, merge the cells vertically
            cell_padding_top = Pt(4)
            if station_name == prev_station_name:
                merged_cell = table.cell(row_idx - 1, num_cols - 1).merge(table.cell(row_idx, num_cols - 1))
                table.cell(row_idx - 1, 0).merge(table.cell(row_idx, 0))

            # Add station name to the appropriate row
            table.cell(row_idx, 0).text = station_name

            # Extract datetime UTC values and values from the data entries
            datetime_values = []
            for entry in data_entries:
                if isinstance(entry, dict) and "dateTimeUTC" in entry:
                    datetime_values.append(entry["dateTimeUTC"])
            # print(type(data_entries))
            # Modify your existing code for extracting values from data_entries
            values = []
            for entry in data_entries:
                if isinstance(entry, dict):
                    if data_type == "rainfall" and "value" in entry:
                        values.append(entry["value"])
                    elif data_type == "cloud_cover" and "oktas" in entry and "cloudgroup" in entry:
                        values.append(f"{entry['oktas']}/{entry['cloudgroup']}")
                    else:
                        values.append("-")  # Handle other cases as needed
            # print(f"Values: {values}")
            # Fill in the table with data using unique datetime values
            for col_idx, datetime_utc in enumerate(unique_datetime_utc_values, start=1):
                value = "-"
                for entry_datetime, entry_value in zip(datetime_values, values):
                    if entry_datetime == datetime_utc:
                        value = entry_value
                        break
                table.cell(row_idx, col_idx).text = value

            # Set the "24-Hour Observed" cell
            value = station_data["total"]
            if value == "T":
                table.cell(row_idx, num_cols - 1).text = "T"
            else:
                # Convert the value to float before formatting
                formatted_value = "{:.1f}".format(float(value))
                table.cell(row_idx, num_cols - 1).text = formatted_value




            # Update the previous station name
            prev_station_name = station_name

            # Replace 'N/A' with '0.0'
            values = [value.replace("N/A", "0.0") for value in values]

            if data_type == "rainfall":
                # Initialize difference variables with default values
                diff_06z_03z = 0.0
                diff_12z_09z = 0.0
                diff_18z_12z = 0.0
                diff_21z_00z = 0.0

                # Check if the values list has enough elements before calculating differences
                if len(values) >= 8:
                    try:
                        # Temporarily replace 'T' with '0.0' for the difference calculation
                        temp_values = [v if v != 'T' else '0.0' for v in values]
                        diff_06z_03z = abs(float(temp_values[1]) - float(temp_values[0]))
                        diff_12z_09z = abs(float(temp_values[3]) - float(temp_values[2]))
                        diff_18z_12z = abs(float(temp_values[5]) - float(temp_values[4]))
                        diff_21z_00z = abs(float(temp_values[7]) - float(temp_values[6]))
                    except ValueError:
                        # Handle the case where a value cannot be converted to float
                        diff_06z_03z = 0.0
                        diff_12z_09z = 0.0
                        diff_18z_12z = 0.0
                        diff_21z_00z = 0.0
                else:
                    # Handle the case where there are not enough values for differences
                    table.cell(row_idx, num_cols - 4).text = "N/A"
                    table.cell(row_idx, num_cols - 3).text = "N/A"
                    table.cell(row_idx, num_cols - 2).text = "N/A"
                    table.cell(row_idx, num_cols - 1).text = "N/A"

                # Update the cells with difference values, retaining 'T' values
                table.cell(row_idx, 2).text = f"{diff_06z_03z:.1f}" if values[1] != 'T' else 'T'
                table.cell(row_idx, 4).text = f"{diff_12z_09z:.1f}" if values[3] != 'T' else 'T'
                table.cell(row_idx, 6).text = f"{diff_18z_12z:.1f}" if values[5] != 'T' else 'T'
                table.cell(row_idx, 8).text = f"{diff_21z_00z:.1f}" if values[7] != 'T' else 'T'



       

            # Update the flag to alternate between rainfall and cloud data
            is_rainfall_data = not is_rainfall_data
            # Add color to columns 2-9 based on the specified value ranges, including "0.0" as gray
            for col_idx in range(1, min(9, num_cols - 1)):  # Adjusted the range for the columns
                try:
                    column_value = float(table.cell(row_idx, col_idx).text_frame.text)
                except (ValueError, AttributeError):
                # Handle the case where the value cannot be converted to a float
                     column_value = 0.0  # Set a default value, or choose another appropriate action
                
                if column_value == 0.0:
                    # 0.0 is gray
                    column_color = RGBColor(255, 255, 255)  # White
                elif column_value < 7.5:
                    # < 7.5 is yellow
                    column_color = RGBColor(255, 255, 0)  # Yellow
                elif 7.5 <= column_value < 22.5:
                    # 7.5-22.5 is orange
                    column_color = RGBColor(255, 165, 0)  # Orange
                elif 22.5 <= column_value < 45:
                    # 22.5-45 is red
                    column_color = RGBColor(255, 0, 0)  # Red
                elif 45 <= column_value < 90:
                    # 45-90 is pink
                    column_color = RGBColor(255, 192, 203)  # Pink
                else:
                    # > 90 is violet
                    column_color = RGBColor(238, 130, 238)  # Violet

                cell = table.cell(row_idx, col_idx)
                fill = cell.fill
                fill.solid()
                fill.fore_color.rgb = column_color


            # Update row index
            row_idx += 1
            # If the maximum rows per slide is reached, create a new slide
            if row_idx > max_rows_per_slide:
                row_idx = 1  # Reset row index

    # Iterate through the slides and apply changes to the 5th slide
    for idx, slide in enumerate(presentation.slides):
        if idx == 4:  # Check if it's the 5th slide (index starts from 0)
            for shape in slide.shapes:
                if not shape.has_table:
                    continue
                table = shape.table

                # Check if the table has at least four rows
                if len(table.rows) > 4:
                    # Iterate over the last four rows of the table
                    for row_idx in range(len(table.rows) - 4, len(table.rows)):
                        row = table.rows[row_idx]

                        # Set the background color of each cell in the row to white
                        for cell in row.cells:
                            cell.text = ""  # Clear cell text
                            fill = cell.fill
                            fill.solid()
                            fill.fore_color.rgb = RGBColor(255, 255, 255)  # White

   


    # Iterate through the slides and add the images to the bottom
    for slide_index, slide in enumerate(presentation.slides):
        # remove the legend in first slide
        if slide_index == 0:
            # Skip adding the legend images for the first slide
            continue
        # Calculate the position and size for the rainfall legend image
        rainfall_legend_left = Inches(0.5)  # Adjust the left position as needed
        rainfall_legend_top = Inches(6.5)   # Adjust the top position as needed
        rainfall_legend_width = Inches(4.5)  # Adjust the width as needed
        rainfall_legend_height = Inches(1.0) # Adjust the height as needed

        # Calculate the position and size for the cloud legend image
        cloud_legend_left = Inches(5.5)  # Adjust the left position as needed
        cloud_legend_top = Inches(6.5)   # Adjust the top position as needed
        cloud_legend_width = Inches(4.5)  # Adjust the width as needed
        cloud_legend_height = Inches(1.0) # Adjust the height as needed

        

        # Add the rainfall legend image to the slide
        add_legend_image(slide, f'{folder}/images/rainfall_legend.png', rainfall_legend_left, rainfall_legend_top, rainfall_legend_width, rainfall_legend_height)


        # Add the cloud legend image to the slide
        add_legend_image(slide, f'{folder}/images/cloud_legend.png', cloud_legend_left, cloud_legend_top, cloud_legend_width, cloud_legend_height)

        table = None
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                # Iterate through all cells in the table
                for row in table.rows:
                    cell.height = Inches(0.5)
                    for cell in row.cells:
                        # Set the desired font size and family for each cell
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(12)  # Set the desired font size
                                run.font.name = "Arial"  # Set the desired font family
                                paragraph.alignment = PP_ALIGN.CENTER  # Set the desired alignment
                                cell.margin_top = Pt(4)
        if table:
            for row in table.rows:
                cell = row.cells[num_cols - 1]  # Get the cell in the last column
                try:
                    column_value = float(cell.text)
                    if column_value == 0.0:
                        # 0.0 is gray
                        fill_color = RGBColor(169, 169, 169)  # Gray
                    elif 0 < column_value < 60:
                        # > 0 and < 60 is yellow
                        fill_color = RGBColor(255, 255, 0)  # Yellow
                    elif 60 <= column_value <= 180:
                        # 60-180 is orange
                        fill_color = RGBColor(255, 165, 0)  # Orange
                    elif 180 < column_value <= 360:
                        # 180-360 is red
                        fill_color = RGBColor(255, 0, 0)  # Red
                    elif 360 < column_value <= 720:
                        # 360-720 is pink
                        fill_color = RGBColor(255, 192, 203)  # Pink
                    elif column_value > 720:
                        # > 720 is violet
                        fill_color = RGBColor(238, 130, 238)  # Violet
                    else:
                        # Default to white for values outside the specified ranges
                        fill_color = RGBColor(255, 255, 255)  # White

                    # Set the cell fill color
                    fill = cell.fill
                    fill.solid()
                    fill.fore_color.rgb = fill_color

                    # Set the font size for the last column
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(14)  # Set the desired font size for the last column
                            run.font.bold = True  # Add bold formatting
                            cell.margin_top = Pt(20)

                except ValueError:
                    pass  # Handle the case where the value is not a float

    return presentation


# Create a PowerPoint presentation
presentation = Presentation()

#value of first image
image_filename = script_name + '.png'

#value of second image
# Convert script_name to a datetime object
script_date = datetime.strptime(script_name, '%Y%m%d')

# Subtract one day
date_today_minus_1 = script_date - timedelta(days=1)
additional_image_filename = date_today_minus_1.strftime('%Y%m%d') 
final = additional_image_filename + '.png'

# print(f"Original script_name: {script_name}")
# print(f"Modified additional_image_filename: {additional_image_filename}")

# Add the default slide with the main image and the additional image on the first slide
add_default_slide(presentation, image_filename, final) 
#add_default_slide(presentation, image_filename)

# Define the path to your local JSON files
rainfall_json_file_path = f'http://10.11.1.107/api/rainfall/{script_name}'
cloud_json_file_path = f'http://10.11.1.107/api/cloud-cover/{script_name}'


# Create slides for Mindanao, Visayas, and Luzon and add them to the single presentation
create_presentation(presentation, "Luzon", rainfall_json_file_path, cloud_json_file_path)
create_presentation(presentation, "Visayas", rainfall_json_file_path, cloud_json_file_path)
create_presentation(presentation, "Mindanao", rainfall_json_file_path, cloud_json_file_path)

# Save the single PowerPoint presentation to a file
presentation.save(f'{folder}/powerpoints/MapDiscussion_{script_name}.pptx')

# COPY TO TERA CODE
transfer = os.path.join(folder, '../copytotera')
sys.path.append(transfer)

from copytotera import mount_network_server, transfer_to_nas, nas_local_dir

# Define the base remote path
base_remote_path = 'wd.s.dstor.pagasa.local/wfs/Rainfall_table_for_Map_Discussion/'

# Get current year and month
current_year = datetime.now().strftime("%Y")
current_year_2_digit = datetime.now().strftime("%y")
current_month = datetime.now().strftime("%m")
month_name = datetime.now().strftime("%b").upper()

# Convert the current month to 3-character format with leading zeros
month_number = datetime.now().strftime("%m")

# Construct the remote path with current year and month
nas_remote_path = f"{base_remote_path}{current_year}/{month_number} {month_name}'{current_year_2_digit}"

# Mount the network server
mount_network_server(nas_remote_path, nas_local_dir)

# Transfer the source image to the NAS local directory
transfer_to_nas(f'{folder}/powerpoints/MapDiscussion_{script_name}.pptx', nas_local_dir)
